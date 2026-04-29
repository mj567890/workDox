"""Text extraction from common document formats.

Supported formats:
  - Binary structured: docx, pdf, xlsx, xls, pptx
  - Plain text (UTF-8/GBK/GB2312/Latin-1): txt, md, csv, json, xml,
    html, htm, log, yaml, yml, toml, cfg, ini, conf, properties, env,
    rst, tex, py, js, ts, tsx, jsx, vue, css, scss, less, sass,
    sql, sh, bash, zsh, bat, cmd, ps1, java, go, rs, c, cpp, cxx,
    h, hpp, hxx, rb, php, swift, kt, scala, r, m, mm, pl, pm, lua,
    dockerfile, makefile, gradle, cmake, gitignore, editorconfig
"""

import io
import traceback
import logging
from pathlib import Path

logger = logging.getLogger(__name__)

# Lazy imports to avoid hard dependency failures at module load time
_import_errors: dict[str, str] = {}


def _get_docx_doc():
    try:
        import docx
        return docx
    except ImportError as e:
        _import_errors["docx"] = str(e)
        return None


def _get_pdf_plumber():
    try:
        import pdfplumber
        return pdfplumber
    except ImportError as e:
        _import_errors["pdfplumber"] = str(e)
        return None


def _get_openpyxl():
    try:
        import openpyxl
        return openpyxl
    except ImportError as e:
        _import_errors["openpyxl"] = str(e)
        return None


def _get_pptx():
    try:
        from pptx import Presentation
        return Presentation
    except ImportError as e:
        _import_errors["pptx"] = str(e)
        return None


MAX_EXTRACT_CHARS = 50000  # Soft cap to avoid storing huge texts

# ── Plain-text file types (use _extract_txt with encoding detection) ──────
TEXT_BASED_TYPES: frozenset[str] = frozenset({
    "txt", "md", "markdown", "rst", "tex",
    "csv", "tsv",
    "json", "jsonc", "jsonl",
    "xml", "svg",
    "html", "htm", "xhtml",
    "yaml", "yml",
    "toml", "cfg", "ini", "conf", "config", "properties", "env",
    "log", "text",
    "py", "js", "ts", "tsx", "jsx", "mjs", "cjs",
    "vue", "svelte",
    "css", "scss", "sass", "less",
    "sql", "psql", "mysql",
    "sh", "bash", "zsh", "fish", "bat", "cmd", "ps1",
    "java", "kt", "kts", "scala", "groovy",
    "go", "rs", "c", "cpp", "cxx", "cc", "h", "hpp", "hxx",
    "rb", "php", "swift", "r", "m", "mm",
    "pl", "pm", "lua", "dart", "elm", "erl", "ex", "exs", "hs", "clj", "cljs",
    "dockerfile", "makefile", "cmake", "gradle",
    "gitignore", "gitattributes", "editorconfig",
    "nginx", "apache",
})

# ── Binary structured types (have dedicated extractors) ────────────────────
BINARY_STRUCTURED_TYPES: frozenset[str] = frozenset({
    "docx", "pdf", "xlsx", "xls", "pptx",
})

# All currently supported extractable types
SUPPORTED_TYPES: frozenset[str] = TEXT_BASED_TYPES | BINARY_STRUCTURED_TYPES


def is_supported(file_type: str) -> bool:
    """Check whether `file_type` has a text extractor."""
    return file_type.lower() in SUPPORTED_TYPES


def _normalize_ext(file_type: str) -> str:
    """Normalize a file-type string to a canonical extension."""
    ft = file_type.lower().strip().lstrip(".")
    # Handle compound extensions like .tar.gz
    if ft.endswith(".gz"):
        ft = ft[:-3]
    if ft.endswith(".tar"):
        ft = ft[:-4]
    return ft


def extract_text(file_data: bytes, file_type: str, original_name: str = "") -> str | None:
    """Extract readable text from a file's binary content.

    Returns extracted text string, or None if extraction fails or is unsupported.
    """
    ft = _normalize_ext(file_type)

    if not ft:
        return None

    try:
        # ── Binary structured formats ──────────────────────────────
        if ft == "docx":
            return _extract_docx(file_data)

        if ft == "pdf":
            return _extract_pdf(file_data)

        if ft in ("xlsx", "xls"):
            return _extract_xlsx(file_data)

        if ft == "pptx":
            return _extract_pptx(file_data)

        # ── Plain-text formats (encoding auto-detect) ──────────────
        if ft in TEXT_BASED_TYPES:
            return _extract_txt(file_data)

        # ── Unsupported ────────────────────────────────────────────
        logger.debug(f"No text extractor for file type: {ft}")
        return None

    except Exception:
        logger.warning(
            f"Text extraction failed for {original_name} ({ft}):\n{traceback.format_exc()}"
        )
        return None


# ═══════════════════════════════════════════════════════════════════════════════
# Individual extractors
# ═══════════════════════════════════════════════════════════════════════════════


def _extract_txt(data: bytes) -> str:
    """Extract text from plain text files with encoding auto-detection."""
    for encoding in ("utf-8", "gbk", "gb2312", "gb18030", "latin-1"):
        try:
            text = data.decode(encoding)
            return text[:MAX_EXTRACT_CHARS]
        except UnicodeDecodeError:
            continue
    return data.decode("utf-8", errors="replace")[:MAX_EXTRACT_CHARS]


def _extract_docx(data: bytes) -> str | None:
    """Extract text from .docx Word documents."""
    docx_mod = _get_docx_doc()
    if docx_mod is None:
        logger.warning("python-docx not installed, skipping docx extraction")
        return None

    doc = docx_mod.Document(io.BytesIO(data))
    paragraphs = []
    for para in doc.paragraphs:
        if para.text.strip():
            paragraphs.append(para.text.strip())

    # Also extract text from tables
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
            if row_text:
                paragraphs.append(row_text)

    text = "\n".join(paragraphs)
    return text[:MAX_EXTRACT_CHARS] if text.strip() else None


def _extract_pdf(data: bytes) -> str | None:
    """Extract text from PDF files using pdfplumber, fallback to pdfminer, then pypdfium2."""
    # Try pdfplumber first (best quality)
    pdfplumber_mod = _get_pdf_plumber()
    if pdfplumber_mod is not None:
        try:
            pages_text = []
            with pdfplumber_mod.open(io.BytesIO(data)) as pdf:
                for page in pdf.pages:
                    page_text = page.extract_text()
                    if page_text:
                        pages_text.append(page_text)
            text = "\n".join(pages_text)
            if text.strip():
                return text[:MAX_EXTRACT_CHARS]
        except Exception:
            logger.debug("pdfplumber extraction failed, trying pdfminer fallback")

    # Fallback to pdfminer.six
    try:
        from pdfminer.high_level import extract_text as pdfminer_extract
        text = pdfminer_extract(io.BytesIO(data))
        if text.strip():
            return text[:MAX_EXTRACT_CHARS]
    except Exception:
        pass

    # Fallback to pypdfium2
    try:
        import pypdfium2
        pdf = pypdfium2.PdfDocument(data)
        pages_text = []
        for page in pdf:
            text_page = page.get_textpage()
            pages_text.append(text_page.get_text_range())
        pdf.close()
        text = "\n".join(pages_text)
        if text.strip():
            return text[:MAX_EXTRACT_CHARS]
    except Exception:
        pass

    return None


def _extract_xlsx(data: bytes) -> str | None:
    """Extract text from Excel spreadsheets."""
    openpyxl_mod = _get_openpyxl()
    if openpyxl_mod is None:
        logger.warning("openpyxl not installed, skipping xlsx extraction")
        return None

    wb = openpyxl_mod.load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    all_text = []

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        sheet_lines = [f"[Sheet: {sheet_name}]"]
        for row in ws.iter_rows(values_only=True):
            row_values = [str(cell) for cell in row if cell is not None]
            if row_values:
                sheet_lines.append(" | ".join(row_values))
        all_text.append("\n".join(sheet_lines))

    wb.close()
    text = "\n\n".join(all_text)
    return text[:MAX_EXTRACT_CHARS] if text.strip() else None


def _extract_pptx(data: bytes) -> str | None:
    """Extract text from .pptx PowerPoint presentations."""
    Presentation = _get_pptx()
    if Presentation is None:
        logger.warning("python-pptx not installed, skipping pptx extraction")
        return None

    prs = Presentation(io.BytesIO(data))
    slides_text = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        slide_parts = [f"[Slide {slide_num}]"]
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    para_text = paragraph.text.strip()
                    if para_text:
                        slide_parts.append(para_text)
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                    if row_text:
                        slide_parts.append(row_text)
        if len(slide_parts) > 1:  # Has more than just the slide header
            slides_text.append("\n".join(slide_parts))

    text = "\n\n".join(slides_text)
    return text[:MAX_EXTRACT_CHARS] if text.strip() else None
